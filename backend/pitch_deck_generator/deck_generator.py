```python
from pptx import Presentation
from pptx.util import Inches, Pt
from backend.database.models.partnership import Partnership

class DeckGenerator:
    def __init__(self):
        self.template_path = "path_to_template.pptx"  # Replace with actual path to template

    def generate_deck(self, partnership_id):
        partnership = self.get_partnership(partnership_id)
        if not partnership:
            return "Partnership not found"

        presentation = Presentation(self.template_path)

        # Customize the slides as per the partnership details
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if "INFLUENCER_NAME" in shape.text:
                        shape.text = shape.text.replace("INFLUENCER_NAME", partnership.influencer.name)
                    if "BRAND_NAME" in shape.text:
                        shape.text = shape.text.replace("BRAND_NAME", partnership.brand.name)
                    if "PARTNERSHIP_IDEA" in shape.text:
                        shape.text = shape.text.replace("PARTNERSHIP_IDEA", partnership.idea)

        # Save the presentation
        presentation.save(f"pitch_deck_{partnership_id}.pptx")

        return f"Pitch deck for partnership {partnership_id} generated successfully"

    @staticmethod
    def get_partnership(partnership_id):
        # Fetch partnership from database
        partnership = Partnership.query.get(partnership_id)
        return partnership
```